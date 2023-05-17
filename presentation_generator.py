import pandas as pd
from Toolkit import Position,add_employees_to_slide
import pptx
import datetime
from data_preparation import data_preparation_processor



def generate_presentation_from_template(prs: pptx.presentation.Presentation, raw_df: pd.DataFrame) -> None:
    # Get separate dataframes for different teams and job families
    teams_df, job_families_df = data_preparation_processor(raw_df)

    business_research_df, office_management_df, business_translation_df, it_df, finance_df,\
        graphic_design_df, marketing_df, operations_df, data_analytics_df, business_development_df,\
        human_capital_df, executive_management_df, sales_representative_df = job_families_df

    mck_team_df, office_admin_team_df, ik_team_df, bcg_me_team_df, business_translation_team_df,\
        it_team_df, finance_team_df, bcg_europe_team_df, graphic_design_team_df, bcg_africa_team_df,\
        marketing_team_df, bcg_kt_support_df, operations_team_df, data_analytics_services_df,\
        business_development_team_df, hr_team_df, research_team_df, executive_management_df, ikt_finance_df = teams_df

    # retrieve all the job families in seperate dfs
    ########## The Cover SLide : Slide[0] , displays the presentation title : 'WHo's WHo ' and the date in ' Month/Year Format #########################################
    # add the current date to the slide
    # get slide 0 of the presentation
    slide = prs.slides[0]
    # get the first textbox shape in the slide
    textbox1 = slide.shapes[0]
    # create a new textbox shape
    textbox2 = slide.shapes.add_textbox(textbox1.left, textbox1.top + textbox1.height + 20, textbox1.width, 50)
    # add the current month and year to the textbox
    today = datetime.date.today()
    month_year_str = today.strftime('%B %Y')
    textbox2.text_frame.text = month_year_str

    # format the text in the new textbox
    font = textbox2.text_frame.paragraphs[0].font
    font.size = pptx.util.Pt(30)  # set font size to 20 points
    font.bold = True  # set bold text
    font.color.rgb = pptx.dml.color.RGBColor(255, 255, 255)  # set white text color

    ############ SLIDE [1] : HARD CODED : Managment Commitee ###############################
    #### This slide is directly derived from the input presentation : it remains unchaged till a further notice ###########

    ######################### SLIDE [2] : HARD CODED :  #######################################################
    #### This slide is directly derived from the input presentation : it remains unchaged till a further notice #################


    ######################### SLIDE [3] :   #######################################################
    # create dfs to be used in slide 3
    bd_officer_df = business_development_df[(business_development_df['Is_officer'] == 1)]
    bd_sr_vp_df = business_development_df[(business_development_df['Is_vice_president'] == 1) & (business_development_df['Is_senior'] == 1)]
    bd_vp_df = business_development_df[(business_development_df['Is_vice_president'] == 1) & (business_development_df['Is_senior'] == 0) & (business_development_df['Is_associate'] == 0)]
    bd_associate_vps_df = business_development_df[(business_development_df['Is_associate'] == 1) & (business_development_df['Is_vice_president'] == 1)]
    bd_lead_df = business_development_df[(business_development_df['Is_lead'] == 1)]
    bd_sr_acc_exec_df = business_development_df[(business_development_df['Is_acc_exec'] == 1) & (business_development_df['Is_senior'] == 1)]
    bd_acc_exec_df = business_development_df[(business_development_df['Is_acc_exec'] == 1) & (business_development_df['Is_senior'] == 0)]

    test_pics_path = './test_data/'

    slide = prs.slides[3]
    add_employees_to_slide(df=bd_officer_df, row_length=2, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_MIDDLE_1)
    add_employees_to_slide(df=bd_sr_vp_df, row_length=3, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_4)
    add_employees_to_slide(df=bd_vp_df, row_length=4, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_MIDDLE_4)
    add_employees_to_slide(df=bd_associate_vps_df, row_length=8, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_1)
    add_employees_to_slide(df=bd_lead_df, row_length=4, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_MIDDLE_1)
    add_employees_to_slide(df=bd_sr_acc_exec_df, row_length=5, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_3)
    add_employees_to_slide(df=bd_acc_exec_df, row_length=5, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_MIDDLE_3)

    ######################### SLIDE [4] :  #######################################################
    test_pics_path = './test_data/'

    # add_employees_to_slide(df= BD_team_df, row_length=3, font_size=7, slide= prs.slides[3] , pics_path = pics_path , detailed_slide=1 , position= Position.TOP_LEFT)
    slide = prs.slides[4]
    # Finance
    add_employees_to_slide(df=finance_team_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)
    # HR
    add_employees_to_slide(df=hr_team_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_4)

    # Operations
    add_employees_to_slide(df=operations_team_df, row_length=8, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_2)
    ######################### SLIDE [5] :   #######################################################

    test_pics_path = './test_data/'
    slide = prs.slides[5]
    # IT
    add_employees_to_slide(df=it_team_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_2)
    # Marketing
    add_employees_to_slide(df=marketing_team_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_LEFT)
    # Office Admin
    add_employees_to_slide(df=office_admin_team_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_2)
    ######################### SLIDE [6] :   #######################################################
    ######## The code commented below works , but this slide better be hard coded as well #######
    test_pics_path = './test_data/'
    finance_directors_df = finance_team_df[finance_team_df['Is_director'] == 1]
    finance_senior_coordinators_df = finance_team_df[
        (finance_team_df['Is_senior'] == 1) & (finance_team_df['Is_coordinator'] == 1)]
    finance_others_df = finance_team_df[~((finance_team_df['Is_director'] == 1) | (
            (finance_team_df['Is_senior'] == 1) & (finance_team_df['Is_coordinator'] == 1)))]

    slide = prs.slides[6]
    # Finance Strategic Topics
    add_employees_to_slide(df=finance_directors_df, row_length=2, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_MIDDLE_3)
    # Business Controller Topics, supporting BD team
    add_employees_to_slide(df=finance_senior_coordinators_df, row_length=2, font_size=8, slide=slide,
                           pics_path=test_pics_path, detailed_slide=1, position=Position.BOTTOM_MIDDLE_1)
    # Country Finance Topics
    add_employees_to_slide(df=finance_others_df, row_length=2, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_RIGHT_2)

    ######################### SLIDE [7] :   #######################################################
    ########## This SLide is left blank : pending detailed Data ###################################

    ######################### SLIDE [8] :   #######################################################
    test_pics_path = './test_data/'
    COO_df = operations_team_df[operations_team_df['Is_COO'] == 1]
    partners_df = operations_team_df[operations_team_df['Is_partner'] == 1]
    operations_team_others_df = operations_team_df[
        (operations_team_df['Is_COO'] != 1) & (operations_team_df['Is_partner'] != 1)]

    slide = prs.slides[8]
    add_employees_to_slide(df=COO_df, row_length=1, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1,
                           position=Position.MIDDLE_LEFT_3)
    add_employees_to_slide(df=operations_team_others_df, row_length=4, font_size=8, slide=slide,
                           pics_path=test_pics_path,
                           detailed_slide=1, position=Position.MAIN_MIDDLE_LEFT)

    ######################### SLIDE [9] :   #######################################################

    test_pics_path = './test_data/'
    slide = prs.slides[9]
    add_employees_to_slide(df=it_team_df, row_length=4, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_LEFT_3)

    ######################### SLIDE [10] :   #######################################################

    test_pics_path = './test_data/'
    slide = prs.slides[10]
    add_employees_to_slide(df=marketing_team_df, row_length=4, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_LEFT_3)

    ######################### SLIDE [11] :   #######################################################

    test_pics_path = './test_data/'
    slide = prs.slides[11]
    add_employees_to_slide(df=office_admin_team_df, row_length=3, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_LEFT_3)
    ######################### SLIDE [12] :   #######################################################
    # create dfs to be used in slide 12
    # Create a dataframe of only the research team head(s)
    research_team_df = business_research_df
    research_head_df = research_team_df[research_team_df['Is_Head'] == 1]
    # Create a dataframe of only the research managers
    research_managers_df = research_team_df[research_team_df['Is_manager'] == 1]
    # Create a dataframe of only the senior research managers
    research_sr_managers_df = research_team_df[
        (research_team_df['Is_manager'] == 1) & (research_team_df['Is_senior'] == 1)]
    # Create a dataframe of only the research team leads
    research_team_leads_df = research_team_df[research_team_df['Is_team_lead'] == 1]

    test_pics_path = './test_data/'
    slide = prs.slides[12]
    add_employees_to_slide(df=research_head_df, row_length=1, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_MIDDLE_2)
    add_employees_to_slide(df=research_sr_managers_df, row_length=4, font_size=6, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_5)
    add_employees_to_slide(df=research_managers_df, row_length=9, font_size=6, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE)
    add_employees_to_slide(df=research_team_leads_df, row_length=10, font_size=6, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_3)
    ######################### SLIDE [13] -1 :   #######################################################
    # create dfs to be used in slide 14
    # Create a dataframe of only the research team head(s)
    research_team_df = business_research_df

    # Create a dataframe of only the senior research associates
    research_sr_associates_df = research_team_df[
        (research_team_df['Is_associate'] == 1) & (research_team_df['Is_senior'] == 1)]

    # Create a dataframe of only the research associates
    research_associates_df = research_team_df[
        (research_team_df['Is_associate'] == 1) & (research_team_df['Is_senior'] == 0)]

    test_pics_path = './test_data/'
    slide = prs.slides[13]
    add_employees_to_slide(df=research_sr_associates_df, row_length=12, font_size=8, slide=slide,
                           pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)
    add_employees_to_slide(df=research_associates_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_1)

    ######################### SLIDE [14] :   #######################################################

    # Create a dataframe of only the senior research senior analysts
    research_sr_analyst_df = research_team_df[
        (research_team_df['Is_analyst'] == 1) & (research_team_df['Is_senior'] == 1)]
    test_pics_path = './test_data/'
    slide = prs.slides[14]
    add_employees_to_slide(df=research_sr_analyst_df, row_length=13, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)

    ######################### SLIDE [15] :   #######################################################
    # create dfs to be used in slide 16
    # Create a dataframe of only the research analysts
    research_analyst_df = research_team_df[(research_team_df['Is_analyst'] == 1) & (research_team_df['Is_senior'] == 0)]
    test_pics_path = './test_data/'
    slide = prs.slides[15]
    add_employees_to_slide(df=research_analyst_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)

    ######################### SLIDE [16] :   #######################################################
    # Hard coded : Left as is for the time being
    ######################### SLIDE [17] :   #######################################################
    # create dfs to be used in slide 17
    # Create a dataframe of only the research analysts
    ikt_sr_managers_df = ik_team_df[(ik_team_df['Is_manager'] == 1) & (ik_team_df['Is_senior'] == 1)]
    ikt_managers_df = ik_team_df[(ik_team_df['Is_manager'] == 1) & (ik_team_df['Is_senior'] == 0)]
    ikt_team_leads_df = ik_team_df[(ik_team_df['Is_team_lead'] == 1)]
    ikt_sr_associates_df = ik_team_df[(ik_team_df['Is_associate'] == 1) & (ik_team_df['Is_senior'] == 1)]
    ikt_associates_df = ik_team_df[(ik_team_df['Is_associate'] == 1) & (ik_team_df['Is_senior'] == 0)]
    test_pics_path = './test_data/'
    slide = prs.slides[17]
    add_employees_to_slide(df=ikt_sr_managers_df, row_length=5, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)
    add_employees_to_slide(df=ikt_managers_df, row_length=8, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_MIDDLE_RIGHT)
    add_employees_to_slide(df=ikt_team_leads_df, row_length=4, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_4)
    add_employees_to_slide(df=ikt_sr_associates_df, row_length=7, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_RIGHT)
    add_employees_to_slide(df=ikt_associates_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_2)

    ######################### SLIDE [18] :   #######################################################
    # create dfs to be used in slide 18
    # Create a dataframe of only the research analysts
    ikt_sr_analysts_df = ik_team_df[(ik_team_df['Is_analyst'] == 1) & (ik_team_df['Is_senior'] == 1)]
    ikt_analysts_df = ik_team_df[(ik_team_df['Is_analyst'] == 1) & (ik_team_df['Is_senior'] == 0)]

    test_pics_path = './test_data/'
    slide = prs.slides[18]
    add_employees_to_slide(df=ikt_sr_analysts_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)
    add_employees_to_slide(df=ikt_analysts_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_1)

    ######################### SLIDE [19] :   #######################################################
    # create dfs to be used in slide 20
    # Create a dataframe of only the research analysts
    mck_sr_managers_df = mck_team_df[(mck_team_df['Is_manager'] == 1) & (mck_team_df['Is_senior'] == 1)]
    mck_managers_df = mck_team_df[(mck_team_df['Is_manager'] == 1) & (mck_team_df['Is_senior'] == 0)]
    mck_team_leads_df = mck_team_df[(mck_team_df['Is_team_lead'] == 1)]
    mck_sr_associates_df = mck_team_df[(mck_team_df['Is_associate'] == 1) & (mck_team_df['Is_senior'] == 1)]
    mck_associates_df = mck_team_df[(mck_team_df['Is_associate'] == 1) & (mck_team_df['Is_senior'] == 0)]
    test_pics_path = './test_data/'
    slide = prs.slides[19]
    add_employees_to_slide(df=mck_sr_managers_df, row_length=5, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_2)
    add_employees_to_slide(df=mck_managers_df, row_length=4, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_MIDDLE_RIGHT)
    add_employees_to_slide(df=mck_team_leads_df, row_length=5, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_LEFT)
    add_employees_to_slide(df=mck_sr_associates_df, row_length=8, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_RIGHT)
    add_employees_to_slide(df=mck_associates_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_2)

    ######################### SLIDE [20] :   #######################################################
    # create dfs to be used in slide 20
    # Create a dataframe of only the research analysts
    mck_sr_analysts_df = mck_team_df[(mck_team_df['Is_analyst'] == 1) & (mck_team_df['Is_senior'] == 1)]
    mck_analysts_df = mck_team_df[(mck_team_df['Is_analyst'] == 1) & (mck_team_df['Is_senior'] == 0)]

    test_pics_path = './test_data/'
    slide = prs.slides[20]
    add_employees_to_slide(df=mck_sr_analysts_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)
    add_employees_to_slide(df=mck_analysts_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_1)

    ######################### SLIDE [21] :   #######################################################
    # create dfs to be used in slide 21

    bcg_team_df = pd.concat([bcg_me_team_df, bcg_europe_team_df, bcg_africa_team_df, bcg_kt_support_df])

    # Create a dataframe of only the research analysts
    bcg_sr_managers_df = bcg_team_df[(bcg_team_df['Is_manager'] == 1) & (bcg_team_df['Is_senior'] == 1)]
    bcg_managers_df = bcg_team_df[(bcg_team_df['Is_manager'] == 1) & (bcg_team_df['Is_senior'] == 0)]
    bcg_team_leads_df = bcg_team_df[(bcg_team_df['Is_team_lead'] == 1)]
    bcg_sr_associates_df = bcg_team_df[(bcg_team_df['Is_associate'] == 1) & (bcg_team_df['Is_senior'] == 1)]
    bcg_associates_df = bcg_team_df[(bcg_team_df['Is_associate'] == 1) & (bcg_team_df['Is_senior'] == 0)]
    test_pics_path = './test_data/'
    slide = prs.slides[21]

    add_employees_to_slide(df=bcg_sr_managers_df, row_length=5, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_2)
    add_employees_to_slide(df=bcg_managers_df, row_length=8, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_MIDDLE_RIGHT)
    add_employees_to_slide(df=bcg_team_leads_df, row_length=5, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_LEFT)
    add_employees_to_slide(df=bcg_sr_associates_df, row_length=7, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_RIGHT)
    add_employees_to_slide(df=bcg_associates_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_2)

    ######################### SLIDE [22] :   #######################################################
    # create dfs to be used in slide 22
    # Create a dataframe of only the research analysts
    bcg_team_df = pd.concat([bcg_me_team_df, bcg_europe_team_df, bcg_africa_team_df, bcg_kt_support_df])
    bcg_sr_analysts_df = bcg_team_df[(bcg_team_df['Is_analyst'] == 1) & (bcg_team_df['Is_senior'] == 1)]
    bcg_analysts_df = bcg_team_df[(bcg_team_df['Is_analyst'] == 1) & (bcg_team_df['Is_senior'] == 0)]

    test_pics_path = './test_data/'
    slide = prs.slides[22]
    add_employees_to_slide(df=bcg_sr_analysts_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_1)
    add_employees_to_slide(df=bcg_analysts_df, row_length=12, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_1)

    ######################### SLIDE [23] :   #######################################################
    # create dfs to be used in slide 23

    business_translation_manager_df = business_translation_df[business_translation_df['Is_manager'] == 1]
    business_translation_others_df = business_translation_df[~((business_translation_df['Is_manager'] == 1))]

    test_pics_path = './test_data/'
    slide = prs.slides[23]
    add_employees_to_slide(df=business_translation_manager_df, row_length=2, font_size=8, slide=slide,
                           pics_path=test_pics_path, detailed_slide=1, position=Position.TOP_MIDDLE_1)

    add_employees_to_slide(df=business_translation_others_df, row_length=8, font_size=8, slide=slide,
                           pics_path=test_pics_path, detailed_slide=1, position=Position.MIDDLE_LEFT_4)
    ######################### SLIDE [24] :   #######################################################
    # create dfs to be used in slide 24

    # Create a dataframe of only the research analysts
    graphic_design_head_df = graphic_design_team_df[(graphic_design_team_df['Is_Head'] == 1)]
    graphic_design_managers_df = graphic_design_team_df[(graphic_design_team_df['Is_manager'] == 1)]
    graphic_design_leads_df = graphic_design_team_df[(graphic_design_team_df['Is_lead'] == 1)]
    graphic_design_sr_df = graphic_design_team_df[(graphic_design_team_df['Is_senior'] == 1)]
    graphic_design_df = graphic_design_team_df[
        (graphic_design_team_df['Is_designer'] == 1) & (graphic_design_team_df['Is_senior'] == 0) & (
                graphic_design_team_df['Is_junior'] == 0)]
    graphic_design_jr_df = graphic_design_team_df[
        (graphic_design_team_df['Is_designer'] == 1) & (graphic_design_team_df['Is_junior'] == 1)]
    graphic_design_coordinator_df = graphic_design_team_df[(graphic_design_team_df['Is_coordinator'] == 1)]

    test_pics_path = './test_data/'
    slide = prs.slides[24]

    add_employees_to_slide(df=graphic_design_head_df, row_length=2, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=1, position=Position.TOP_LEFT_1)
    add_employees_to_slide(df=graphic_design_managers_df, row_length=4, font_size=8, slide=slide,
                           pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_MIDDLE_RIGHT)
    add_employees_to_slide(df=graphic_design_leads_df, row_length=3, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_RIGHT_1)

    add_employees_to_slide(df=graphic_design_sr_df, row_length=3, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.TOP_LEFT_5)

    add_employees_to_slide(df=graphic_design_df, row_length=4, font_size=7, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.MAIN_MIDDLE_1)

    add_employees_to_slide(df=graphic_design_coordinator_df, row_length=2, font_size=7, slide=slide,
                           pics_path=test_pics_path, detailed_slide=0, position=Position.TOP_RIGHT_5)

    add_employees_to_slide(df=graphic_design_jr_df, row_length=10, font_size=8, slide=slide, pics_path=test_pics_path,
                           detailed_slide=0, position=Position.BOTTOM_LEFT_2)

    ######################### SLIDE [25] :   #######################################################
    # create dfs to be used in slide 25

    data_analytics_services_head_df = data_analytics_services_df[data_analytics_services_df['Is_Head'] == 1]
    data_analytics_services_head_others_df = data_analytics_services_df[~((data_analytics_services_df['Is_Head'] == 1))]

    test_pics_path = './test_data/'
    slide = prs.slides[25]
    add_employees_to_slide(df=data_analytics_services_head_df, row_length=1, font_size=8, slide=slide,
                           pics_path=test_pics_path, detailed_slide=1, position=Position.TOP_MIDDLE_RIGHT)

    add_employees_to_slide(df=data_analytics_services_head_others_df, row_length=5, font_size=7, slide=slide,
                           pics_path=test_pics_path, detailed_slide=1, position=Position.MAIN_MIDDLE_1)

    return prs





