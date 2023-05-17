import os
from enum import Enum
from typing import Tuple
import pptx.util


class Position(Enum):
    """Enumeration of positions to be used when positioning employees on the PowerPoint slide."""
    TOP_LEFT_1 = 'top left 1'
    TOP_LEFT_2 = 'top left 2'
    TOP_LEFT_3 = 'top left 3'
    TOP_LEFT_4 = 'top left 4'
    TOP_LEFT_5 = 'top left 5'
    MIDDLE_LEFT_1 = 'middle left 1'
    MIDDLE_LEFT_2 = 'middle left 2'
    MIDDLE_LEFT_3 = 'middle left 3'
    MIDDLE_LEFT_4 = 'middle left 4'
    TOP_RIGHT_1 = 'top right 1'
    TOP_RIGHT_2 = 'top right 2'
    TOP_RIGHT_3 = 'top right 3'
    TOP_RIGHT_4 = 'top right 4'
    TOP_RIGHT_5 = 'top right 5'
    TOP_MIDDLE_RIGHT ='top middle right'
    MAIN_MIDDLE_LEFT = 'middle left'
    MAIN_MIDDLE_RIGHT = 'main middle right'
    BOTTOM_LEFT_1 = 'bottom left 1'
    BOTTOM_LEFT_2 = 'bottom left 2'
    BOTTOM_LEFT_3 = 'bottom left 3'
    BOTTOM_RIGHT_1 = 'bottom right 1'
    BOTTOM_RIGHT_2 = 'bottom right 2'
    TOP_MIDDLE_1 = 'top middle 1'
    TOP_MIDDLE_2 = 'top middle 2'
    TOP_MIDDLE_3 = 'top middle 3'
    TOP_MIDDLE_4 = 'top middle 4'
    MAIN_MIDDLE = 'main middle'
    MAIN_MIDDLE_1 = 'main middle 1'
    MAIN_MIDDLE_2 = 'main middle 2'
    BOTTOM_MIDDLE_1 = 'bottom middle 1'
    BOTTOM_MIDDLE_2 = 'bottom middle 2'
    BOTTOM_MIDDLE_3 = 'bottom middle 3'


def cursor_positioner(position: Position, col: int, row: int) -> Tuple[pptx.util.Length, pptx.util.Length]:
    """
    Calculate the left and top position of the cursor in a PowerPoint slide.

    Args:
        position (Position): The position of the cursor in the slide.
        col (int): The column of the cursor.
        row (int): The row of the cursor.

    Returns:
        tuple: A tuple containing the left and top position of the cursor.

    Raises:
        AssertionError: If the column and row are not positive integers.
        ValueError: If the position is not valid.
    """
    assert col >= 0, "Column must be a non-negative integer."
    assert row >= 0, "Row must be a non-negative integer."

    position_map = {
        Position.TOP_LEFT_1: (0.5, 1.5),
        Position.TOP_LEFT_2: (0.5, 1.90),
        Position.TOP_LEFT_3: (0.5 , 2.2),
        Position.TOP_LEFT_4: (0.5, 3),
        Position.TOP_LEFT_5: (0.5, 3.2),
        Position.MIDDLE_LEFT_1: (1.5, 1.5),
        Position.MIDDLE_LEFT_2: (1.5, 1.90),
        Position.MIDDLE_LEFT_3: (1.5, 2.2),
        Position.MIDDLE_LEFT_4: (1.5, 3),
        Position.TOP_RIGHT_1: (11, 1.5),
        Position.TOP_RIGHT_2: (11, 1.90),
        Position.TOP_RIGHT_3: (11, 2.2),
        Position.TOP_RIGHT_4: (11, 3),
        Position.TOP_RIGHT_5: (11, 3.2),
        Position.MAIN_MIDDLE_LEFT: (0.5, 4),
        Position.TOP_MIDDLE_RIGHT: (5.5, 1.5),
        Position.MAIN_MIDDLE_RIGHT: (6, 3.2),
        Position.BOTTOM_LEFT_1: (0.5, 4.7),
        Position.BOTTOM_LEFT_2: (0.5, 6),
        Position.BOTTOM_LEFT_3: (0.5, 6.2),
        Position.BOTTOM_RIGHT_1: (12, 4.9),
        Position.BOTTOM_RIGHT_2: (12, 6),
        Position.TOP_MIDDLE_1: (4, 1.5),
        Position.TOP_MIDDLE_2: (4, 1.8),
        Position.TOP_MIDDLE_3: (4, 2.2),
        Position.TOP_MIDDLE_4: (4.2, 3),
        Position.MAIN_MIDDLE: (4.6, 3.2),
        Position.MAIN_MIDDLE_1: (5, 3.5),
        Position.MAIN_MIDDLE_2: (5, 4.5),
        Position.BOTTOM_MIDDLE_1: (6, 4.9),
        Position.BOTTOM_MIDDLE_2: (5, 6),
        Position.BOTTOM_MIDDLE_3: (5, 6.2)
    }

    if position in position_map:
        left_offset, top_offset = position_map[position]
        left = pptx.util.Inches(left_offset + col * 1.0)
        top = pptx.util.Inches(top_offset + row * 1.5)
    else:
        raise ValueError("Invalid position.")

    return left, top


def add_employees_to_slide(df, row_length: int, font_size: int, slide, pics_path: str, detailed_slide: bool,
                           position: str):
    """
    Add employees to a PowerPoint slide.

    Args:
        df (pandas.DataFrame): A DataFrame containing employee data.
        row_length (int): The number of employees per row.
        font_size (int): The font size to use for the employee name and location.
        slide: The slide to add the employees to.
        pics_path (str): The path to the directory containing employee images.
        detailed_slide (bool): Whether to include the job position in the employee details.
        position (str): The position to place the employee images on the slide.

    Raises:
        TypeError: If the row_length or font_size is not an integer.
    """
    if not isinstance(row_length, int):
        raise TypeError("Row length must be an integer.")
    if not isinstance(font_size, int):
        raise TypeError("Font size must be an integer.")

    df.reset_index(inplace=True)
    # Adjust the photo size to the font size
    photo_size = font_size / 10

    for i, row in df.iterrows():
        # Checking and formating full name to fit the textbox
        full_name = row['full_name']
        if ' ' in full_name:
            name_parts = full_name.split(' ')
            first_name_length = len(name_parts[0])
            split_index = next(
                (i for i, _ in enumerate(name_parts) if sum(len(p) for p in name_parts[:i + 1]) > len(full_name) / 2),
                len(name_parts))
            full_name = ' '.join(name_parts[:split_index]) + ' \n' + ' '.join(name_parts[split_index:])

        # Checking and formating job title to fit the textbox
        job_title = row['job_position']
        if len(job_title) - job_title.count(' ') > 12:
            job_parts = job_title.split(' ')
            if len(job_parts) > 1:
                job_title = '\n' + job_parts[0] + '\n' + ' '.join(job_parts[1:])  if detailed_slide else ''
            else:
                job_title = '\n' + job_title if detailed_slide else ''
        else:
            job_title = '\n' + job_title if detailed_slide else ''

        location = row['location']
        empl_image_path = 'test_pic.png'
        image_path = os.path.join(pics_path, empl_image_path)

        col = i % row_length
        row_num = i // row_length
        left, top = cursor_positioner(position, col, row_num)

        picture_shape = slide.shapes.add_picture(
            image_path, left=left, top=top,
            width=pptx.util.Inches(photo_size), height=pptx.util.Inches(photo_size))
        textbox_width = (picture_shape.width)//2
        textbox_height = pptx.util.Inches(0.5)
        # position the textbox underneath the picture , and resize it to fit the text
        textbox = slide.shapes.add_textbox(
            left=left, top=top + picture_shape.height // 1.5,
            width=textbox_width, height=textbox_height)

        textframe = textbox.text_frame
        paragraph1 = textframe.add_paragraph()
        paragraph1.text = f"{full_name}{job_title}\n{location}"
        paragraph1.font.size = pptx.util.Pt(font_size)
